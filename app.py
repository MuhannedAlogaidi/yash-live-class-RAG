import streamlit as st
import os
import tempfile
from pathlib import Path

# Document loaders
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader
from langchain_community.vectorstores import FAISS

# Text splitter (moved to own package in LangChain 1.0)
from langchain_text_splitters import RecursiveCharacterTextSplitter

# Embeddings
from langchain_huggingface import HuggingFaceEmbeddings

# LLM
from langchain_groq import ChatGroq

# Pure langchain_core — no deprecated chains, no langchain.chains at all
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.documents import Document

# PPTX custom loader
from pptx import Presentation

# ─────────────────────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="DocChat",
    page_icon="",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ─── Cached Embedding Model ───────────────────────────────────────────────────
@st.cache_resource(show_spinner="Loading sentence-transformer (first run only)...")
def load_embedding_model() -> HuggingFaceEmbeddings:
    return HuggingFaceEmbeddings(
        model_name="sentence-transformers/all-MiniLM-L6-v2",
        model_kwargs={"device": "cpu"},
        encode_kwargs={"normalize_embeddings": True},
    )


# ─── Document Loaders ─────────────────────────────────────────────────────────

def _load_pdf(path: str) -> list[Document]:
    return PyPDFLoader(path).load()


def _load_docx(path: str) -> list[Document]:
    return Docx2txtLoader(path).load()


def _load_pptx(path: str) -> list[Document]:
    prs = Presentation(path)
    docs = []
    for idx, slide in enumerate(prs.slides, start=1):
        text = "\n".join(
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        )
        if text:
            docs.append(Document(
                page_content=text,
                metadata={"source": path, "slide": idx},
            ))
    return docs


def ingest_uploaded_file(uploaded_file) -> list[Document]:
    suffix = Path(uploaded_file.name).suffix.lower()
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(uploaded_file.getvalue())
        path = tmp.name
    try:
        if suffix == ".pdf":
            docs = _load_pdf(path)
        elif suffix in (".docx", ".doc"):
            docs = _load_docx(path)
        elif suffix in (".pptx", ".ppt"):
            docs = _load_pptx(path)
        else:
            raise ValueError(f"Unsupported file type: {suffix}")
    finally:
        os.unlink(path)
    return docs


# ─── RAG Components ───────────────────────────────────────────────────────────
# Plain LCEL — no langchain.chains, no deprecated imports, no conversation history.
# Each question is answered independently: embed → retrieve → prompt → LLM.

_QA_PROMPT = ChatPromptTemplate.from_template(
    "You are a helpful assistant answering questions about an uploaded document.\n"
    "Use only the context below to answer. If the answer is not in the context, say so.\n\n"
    "Context:\n{context}\n\n"
    "Question: {question}\n\n"
    "Answer:"
)


def build_rag_components(docs: list[Document], groq_api_key: str, model_name: str):
    embeddings = load_embedding_model()

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=50,
        separators=["\n\n", "\n", ".", " ", ""],
    )
    chunks = splitter.split_documents(docs)
    vectorstore = FAISS.from_documents(chunks, embeddings)
    retriever = vectorstore.as_retriever(search_kwargs={"k": 4})

    llm = ChatGroq(
        groq_api_key=groq_api_key,
        model_name=model_name,
        temperature=0.2,
        max_tokens=1024,
    )

    # LCEL chain: fills prompt then calls LLM, returns plain string
    qa_chain = _QA_PROMPT | llm | StrOutputParser()

    return retriever, qa_chain, len(chunks)


def answer_question(retriever, qa_chain, question: str):
    source_docs = retriever.invoke(question)
    context = "\n\n".join(doc.page_content for doc in source_docs)
    answer = qa_chain.invoke({"context": context, "question": question})
    return answer, source_docs


# ─── Architecture Page ────────────────────────────────────────────────────────

def render_stage(title: str, body: str):
    st.markdown(f"**{title}**")
    st.caption(body)


def page_architecture():
    st.title("Pipeline Architecture")
    st.caption(
        "A visual walkthrough of the two pipelines that power this app: "
        "data ingestion (upload) and query (chat)."
    )

    # ── 1. Data Ingestion ──────────────────────────────────────────────────
    st.markdown("### 1  Data Ingestion Pipeline")
    st.caption("Runs once when you upload a document.")

    i_cols = st.columns([3, 1, 3, 1, 3, 1, 3, 1, 3])
    ingestion = [
        ("Upload",   "PDF / DOCX / PPTX received via Streamlit file uploader."),
        ("Extract",  "Text pulled page-by-page (PDF), paragraph-by-paragraph (DOCX), or slide-by-slide (PPTX)."),
        ("Chunk",    "RecursiveCharacterTextSplitter cuts text into 500-token windows with 50-token overlap."),
        ("Embed",    "all-MiniLM-L6-v2 maps each chunk to a 384-dim dense vector (CPU, cached on VM)."),
        ("Index",    "Vectors inserted into an in-memory FAISS flat index, ready for similarity search."),
    ]
    for col, (title, desc) in zip(i_cols[::2], ingestion):
        with col:
            render_stage(title, desc)
    for col in i_cols[1::2]:
        with col:
            st.write("→")

    st.divider()

    # ── 2. Query Pipeline ─────────────────────────────────────────────────
    st.markdown("### 2  Query Pipeline")
    st.caption("Runs on every message you send.")

    q_cols = st.columns([3, 1, 3, 1, 3, 1, 3, 1, 3])
    query = [
        ("Question", "User types a question in the chat input."),
        ("Retrieve",  "FAISS returns top-4 chunks closest to the question vector."),
        ("Prompt",    "Chunks joined as context and inserted into a ChatPromptTemplate."),
        ("Groq LLM",  "LLaMA / Mixtral / Gemma on Groq Cloud generates the grounded answer."),
        ("Reply",     "Answer shown in chat. Source chunks in the expander below the reply."),
    ]
    for col, (title, desc) in zip(q_cols[::2], query):
        with col:
            render_stage(title, desc)
    for col in q_cols[1::2]:
        with col:
            st.write("→")

    st.divider()

    # ── 3. Component Details ──────────────────────────────────────────────
    st.markdown("### 3  Component Details")

    c1, c2, c3 = st.columns(3)

    with c1:
        st.markdown("**Embedding Model**")
        st.code(
            "sentence-transformers/all-MiniLM-L6-v2\n"
            "Package    : langchain-huggingface\n"
            "Dimensions : 384\n"
            "Device     : CPU\n"
            "Caching    : @st.cache_resource\n"
            "             downloaded once, lives in\n"
            "             ~/.cache/huggingface on VM",
            language=None,
        )

        st.markdown("**Vector Store**")
        st.code(
            "Package : langchain-community + faiss-cpu\n"
            "Index   : Flat L2\n"
            "Search  : k-NN, k=4\n"
            "Storage : in-memory (per upload)",
            language=None,
        )

    with c2:
        st.markdown("**Text Splitter**")
        st.code(
            "Package    : langchain-text-splitters\n"
            "Type       : RecursiveCharacterTextSplitter\n"
            "Chunk size : 500 tokens\n"
            "Overlap    : 50 tokens\n"
            "Separators : \\n\\n → \\n → . → ' '",
            language=None,
        )

        st.markdown("**RAG Chain (LCEL)**")
        st.code(
            "ChatPromptTemplate   (langchain_core)\n"
            "  | ChatGroq         (langchain_groq)\n"
            "  | StrOutputParser  (langchain_core)\n\n"
            "No deprecated chains.\n"
            "No conversation history.\n"
            "Each question answered independently.",
            language=None,
        )

    with c3:
        st.markdown("**LLM  (BYOK via Groq)**")
        st.code(
            "Package     : langchain-groq\n"
            "Provider    : Groq Cloud API\n"
            "Auth        : user API key (never stored)\n"
            "Models      :\n"
            "  llama-3.3-70b-versatile\n"
            "  llama-3.1-8b-instant\n"
            "  mixtral-8x7b-32768\n"
            "  gemma2-9b-it\n"
            "Temperature : 0.2\n"
            "Max tokens  : 1024",
            language=None,
        )

    st.divider()

    # ── 4. File Format Support ────────────────────────────────────────────
    st.markdown("### 4  Supported File Formats")
    f1, f2, f3 = st.columns(3)
    with f1:
        st.markdown("**PDF**")
        st.caption("PyPDFLoader (pypdf). Text extracted per page with page-number metadata.")
    with f2:
        st.markdown("**DOCX**")
        st.caption("Docx2txtLoader (docx2txt). Full document text extracted in one pass.")
    with f3:
        st.markdown("**PPTX**")
        st.caption("Custom loader (python-pptx). Text per slide with slide-number metadata. No unstructured dependency.")


# ─── Chat Page ────────────────────────────────────────────────────────────────

def page_chat(uploaded_file, groq_api_key: str, model_name: str):
    st.title("Chat with Your Document")

    for key, val in {
        "messages": [],
        "retriever": None,
        "qa_chain": None,
        "loaded_file": None,
        "stats": None,
    }.items():
        if key not in st.session_state:
            st.session_state[key] = val

    # ── Process new / changed upload ───────────────────────────────────────
    if uploaded_file is not None:
        if not groq_api_key:
            st.warning("Enter your Groq API key in the sidebar to continue.")
            return

        if st.session_state.loaded_file != uploaded_file.name:
            with st.status("Processing document...", expanded=True) as status:
                st.write("Warming up embedding model (cached after first run)...")
                load_embedding_model()

                st.write(f"Extracting text from **{uploaded_file.name}** ...")
                docs = ingest_uploaded_file(uploaded_file)
                st.write(f"Extracted {len(docs)} page(s) / slide(s).")

                st.write("Chunking, embedding, and building FAISS index...")
                retriever, qa_chain, n_chunks = build_rag_components(docs, groq_api_key, model_name)
                st.write(f"Indexed {n_chunks} chunks. Ready.")

                st.session_state.retriever = retriever
                st.session_state.qa_chain = qa_chain
                st.session_state.loaded_file = uploaded_file.name
                st.session_state.messages = []
                st.session_state.stats = {
                    "file": uploaded_file.name,
                    "pages": len(docs),
                    "chunks": n_chunks,
                }
                status.update(label="Ready — start chatting below.", state="complete", expanded=False)

    # ── Stats bar ──────────────────────────────────────────────────────────
    if st.session_state.stats:
        s = st.session_state.stats
        m1, m2, m3 = st.columns(3)
        m1.metric("Document", s["file"])
        m2.metric("Pages / Slides", s["pages"])
        m3.metric("Chunks in FAISS", s["chunks"])
        st.divider()

    # ── Empty state ────────────────────────────────────────────────────────
    if st.session_state.retriever is None:
        st.info(
            "Upload a PDF, DOCX, or PPTX in the sidebar and enter your Groq API key "
            "to start chatting with your document."
        )
        return

    # ── Chat history (display only) ────────────────────────────────────────
    for msg in st.session_state.messages:
        with st.chat_message(msg["role"]):
            st.write(msg["content"])

    # ── Chat input ─────────────────────────────────────────────────────────
    if user_input := st.chat_input("Ask something about your document..."):
        st.session_state.messages.append({"role": "user", "content": user_input})
        with st.chat_message("user"):
            st.write(user_input)

        with st.chat_message("assistant"):
            with st.spinner("Retrieving and generating..."):
                answer, sources = answer_question(
                    st.session_state.retriever,
                    st.session_state.qa_chain,
                    user_input,
                )

            st.write(answer)

            if sources:
                with st.expander("Source chunks used", expanded=False):
                    for i, doc in enumerate(sources[:4], 1):
                        meta = doc.metadata
                        label = f"Chunk {i}"
                        if "page" in meta:
                            label += f"  —  Page {int(meta['page']) + 1}"
                        elif "slide" in meta:
                            label += f"  —  Slide {meta['slide']}"
                        st.caption(label)
                        preview = doc.page_content
                        if len(preview) > 350:
                            preview = preview[:350] + "..."
                        st.text(preview)
                        if i < min(4, len(sources)):
                            st.markdown("---")

        st.session_state.messages.append({"role": "assistant", "content": answer})


# ─── Sidebar + Entry Point ────────────────────────────────────────────────────

def main():
    with st.sidebar:
        st.title("DocChat")
        st.caption("RAG over your own documents")

        st.divider()

        st.subheader("Groq API Key")
        groq_api_key = st.text_input(
            "Paste your key here",
            type="password",
            placeholder="gsk_...",
            help="Free key at console.groq.com — never stored or logged.",
            label_visibility="collapsed",
        )
        if groq_api_key:
            st.caption("Key received. Used only for this session.")
        else:
            st.caption("Get a free key at [console.groq.com](https://console.groq.com)")

        st.divider()

        st.subheader("Model")
        model_name = st.selectbox(
            "model",
            options=[
                "llama-3.3-70b-versatile",
                "llama-3.1-8b-instant",
                "mixtral-8x7b-32768",
                "gemma2-9b-it",
            ],
            label_visibility="collapsed",
        )

        st.divider()

        st.subheader("Navigation")
        page = st.radio(
            "nav",
            options=["Chat", "Architecture"],
            label_visibility="collapsed",
        )

        st.divider()

        uploaded_file = None
        if page == "Chat":
            st.subheader("Upload Document")
            uploaded_file = st.file_uploader(
                "file",
                type=["pdf", "docx", "pptx"],
                help="PDF, Word (.docx), or PowerPoint (.pptx)",
                label_visibility="collapsed",
            )

        st.divider()
        st.caption("LangChain · FAISS · HuggingFace · Groq")

    if page == "Architecture":
        page_architecture()
    else:
        page_chat(uploaded_file, groq_api_key, model_name)


if __name__ == "__main__":
    main()
